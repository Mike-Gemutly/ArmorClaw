"""Shared pytest fixtures for ArmorClaw Python Office Sidecar tests.

Generates real test fixture files for all 6 supported document formats
(.xlsx, .pptx, .msg, .doc, .xls, .ppt) and provides gRPC server/channel
fixtures for integration testing.
"""

import hashlib
import hmac
import io
import os
import struct
import sys
import threading
import time
from concurrent import futures

import grpc
import openpyxl
import pytest
import xlwt
from pptx import Presentation

sys.path.insert(0, os.path.dirname(__file__))

from proto import sidecar_pb2, sidecar_pb2_grpc
from worker import (
    FORMAT_MAP,
    MAX_REQUESTS,
    OfficeSidecarServicer,
    SERVER_VERSION,
    _THRESHOLD_BYTES,
)
from interceptor import (
    MAX_TIMESTAMP_AGE_SECONDS,
    TOKEN_METADATA_KEY,
    TOKEN_TTL_SECONDS,
    SERVER_VERSION_METADATA_KEY,
    _calculate_hmac,
    validate_token,
)


TEST_SECRET = "test-hmac-secret-for-unit-tests"

OLE_MAGIC = b"\xd0\xcf\x11\xe0\xa1\xb1\x1a\xe1"
ZIP_MAGIC = b"PK\x03\x04"



def make_token(
    request_id: str,
    timestamp: int,
    operation: str,
    secret: str = TEST_SECRET,
) -> str:
    """Build a valid HMAC token matching the Go/Python token format."""
    data = f"{request_id}{timestamp}{operation}"
    sig = _calculate_hmac(secret, data)
    return f"{request_id}:{timestamp}:{operation}:{sig}"


def make_token_metadata(
    request_id: str = "test-req",
    timestamp: int | None = None,
    operation: str = "extract_text",
    secret: str = TEST_SECRET,
) -> list[tuple[str, str]]:
    """Return gRPC metadata list with a valid token."""
    if timestamp is None:
        timestamp = int(time.time())
    token = make_token(request_id, timestamp, operation, secret)
    return [(TOKEN_METADATA_KEY, token)]




def _create_minimal_ole(stream_name: str, stream_data: bytes) -> bytes:
    """Create a minimal OLE2 / CFBF file with one named stream.

    The file has exactly 4 sectors of 512 bytes:
      Sector 0 — OLE header (512 bytes)
      Sector 1 — FAT (512 bytes)
      Sector 2 — Directory (512 bytes)
      Sector 3 — Stream data (512 bytes)
    """
    SECTOR = 512
    ENDOFCHAIN = 0xFFFFFFFE
    FATSECT = 0xFFFFFFFD
    FREESECT = 0xFFFFFFFF

    # --- Header (sector 0) ---
    hdr = bytearray(SECTOR)
    hdr[0:8] = OLE_MAGIC
    struct.pack_into("<H", hdr, 0x18, 0x003E)
    struct.pack_into("<H", hdr, 0x1A, 0x0003)  # major version (v3)
    struct.pack_into("<H", hdr, 0x1C, 0xFFFE)
    struct.pack_into("<H", hdr, 0x1E, 9)
    struct.pack_into("<H", hdr, 0x20, 6)
    struct.pack_into("<I", hdr, 0x2C, 1)
    struct.pack_into("<I", hdr, 0x30, 1)
    struct.pack_into("<I", hdr, 0x38, 0x1000)
    struct.pack_into("<I", hdr, 0x3C, ENDOFCHAIN)
    struct.pack_into("<I", hdr, 0x40, 0)
    struct.pack_into("<I", hdr, 0x44, ENDOFCHAIN)
    struct.pack_into("<I", hdr, 0x48, 0)
    # DIFAT[0] = sector 0 (the FAT itself)
    struct.pack_into("<I", hdr, 0x4C, 0)
    for i in range(1, 109):
        struct.pack_into("<I", hdr, 0x4C + i * 4, FREESECT)

    # --- FAT (sector 1) ---
    fat = bytearray(SECTOR)
    struct.pack_into("<I", fat, 0, FATSECT)
    struct.pack_into("<I", fat, 4, FATSECT)
    struct.pack_into("<I", fat, 8, ENDOFCHAIN)
    for i in range(3, 128):
        struct.pack_into("<I", fat, i * 4, FREESECT)

    # --- Directory (sector 2) ---
    direntries = bytearray(SECTOR)

    # Entry 0 — Root Entry
    e0 = bytearray(128)
    name0 = "Root Entry".encode("utf-16-le")
    e0[0 : len(name0)] = name0
    struct.pack_into("<H", e0, 0x40, len(name0) + 2)
    e0[0x42] = 5
    e0[0x43] = 1
    struct.pack_into("<H", e0, 0x44, 1)  # child DID → entry 1

    # Entry 1 — Named stream
    e1 = bytearray(128)
    sname = stream_name.encode("utf-16-le")
    e1[0 : len(sname)] = sname
    struct.pack_into("<H", e1, 0x40, len(sname) + 2)
    e1[0x42] = 2
    e1[0x43] = 1
    struct.pack_into("<I", e1, 0x74, 2)
    struct.pack_into("<I", e1, 0x78, len(stream_data))

    direntries[0:128] = e0
    direntries[128:256] = e1

    # --- Data (sector 3) ---
    data_sector = bytearray(SECTOR)
    data_sector[0 : len(stream_data)] = stream_data

    return bytes(hdr) + bytes(fat) + bytes(direntries) + bytes(data_sector)




class _SyncTokenInterceptor(grpc.ServerInterceptor):
    """Synchronous gRPC server interceptor that validates HMAC tokens.

    The production TokenInterceptor is async (grpc.aio.ServerInterceptor)
    which is incompatible with the synchronous grpc.server() used in tests.
    This sync version reuses the same validate_token() logic.
    """

    def __init__(self, shared_secret: str):
        self._secret = shared_secret

    def intercept_service(self, continuation, handler_call_details):
        metadata = dict(handler_call_details.invocation_metadata or [])
        token = metadata.get(TOKEN_METADATA_KEY)

        if not token:
            msg = "missing token"
            return grpc.unary_unary_rpc_method_handler(
                lambda req, ctx, m=msg: ctx.abort(
                    grpc.StatusCode.UNAUTHENTICATED, m
                )
            )

        try:
            validate_token(self._secret, token)
        except ValueError as e:
            msg = str(e)
            return grpc.unary_unary_rpc_method_handler(
                lambda req, ctx, m=msg: ctx.abort(
                    grpc.StatusCode.UNAUTHENTICATED, m
                )
            )

        return continuation(handler_call_details)




_GRPC_OPTIONS = [
    ("grpc.max_receive_message_length", 100 * 1024 * 1024),
    ("grpc.max_send_message_length", 100 * 1024 * 1024),
]


def _start_grpc_server(
    socket_path: str,
    interceptors: list | None = None,
) -> grpc.Server:
    server = grpc.server(
        futures.ThreadPoolExecutor(max_workers=4),
        interceptors=interceptors or [],
        options=[
            ("grpc.max_receive_message_length", 100 * 1024 * 1024),
            ("grpc.max_send_message_length", 100 * 1024 * 1024),
        ],
    )
    sidecar_pb2_grpc.add_SidecarServiceServicer_to_server(
        OfficeSidecarServicer(), server
    )
    os.umask(0o077)
    server.add_insecure_port(f"unix://{socket_path}")
    server.start()
    return server




@pytest.fixture(scope="session")
def xlsx_file(tmp_path_factory):
    """Minimal .xlsx with 'Hello ArmorClaw' in A1."""
    path = tmp_path_factory.mktemp("fixtures") / "test.xlsx"
    wb = openpyxl.Workbook()
    ws = wb.active
    ws.title = "Sheet"
    ws["A1"] = "Hello ArmorClaw"
    ws["B1"] = "Test Value"
    ws["A2"] = "Second row"
    wb.save(str(path))
    return path


@pytest.fixture(scope="session")
def pptx_file(tmp_path_factory):
    """Minimal .pptx with 1 slide containing 'Test Presentation'."""
    path = tmp_path_factory.mktemp("fixtures") / "test.pptx"
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Test Presentation"
    prs.save(str(path))
    return path


@pytest.fixture(scope="session")
def xls_file(tmp_path_factory):
    """Minimal .xls with 1 sheet, cell A1 = 'Test'."""
    path = tmp_path_factory.mktemp("fixtures") / "test.xls"
    wb = xlwt.Workbook()
    ws = wb.add_sheet("Sheet1")
    ws.write(0, 0, "Test")
    wb.save(str(path))
    return path


@pytest.fixture(scope="session")
def msg_file(tmp_path_factory):
    """Minimal .msg (OLE) with email content."""
    path = tmp_path_factory.mktemp("fixtures") / "test.msg"
    content = _create_minimal_ole(
        "Email",
        b"Subject: Test Email\r\nFrom: test@armorclaw.com\r\nHello from ArmorClaw",
    )
    path.write_bytes(content)
    return path


@pytest.fixture(scope="session")
def doc_file(tmp_path_factory):
    """Minimal .doc (OLE) with WordDocument stream."""
    path = tmp_path_factory.mktemp("fixtures") / "test.doc"
    content = _create_minimal_ole(
        "WordDocument",
        b"Test Document Content\r\nHello ArmorClaw",
    )
    path.write_bytes(content)
    return path


@pytest.fixture(scope="session")
def ppt_file(tmp_path_factory):
    """Minimal .ppt (OLE) with PowerPoint Document stream."""
    path = tmp_path_factory.mktemp("fixtures") / "test.ppt"
    content = _create_minimal_ole(
        "PowerPoint Document",
        b"Test Presentation Content\r\nSlide 1",
    )
    path.write_bytes(content)
    return path


@pytest.fixture(scope="session")
def large_xlsx_file(tmp_path_factory):
    """Large .xlsx exceeding 10MB threshold for streaming tests."""
    import random
    path = tmp_path_factory.mktemp("fixtures") / "large.xlsx"
    wb = openpyxl.Workbook()
    ws = wb.active
    ws.title = "BigSheet"
    rng = random.Random(42)
    row = 0
    while True:
        row += 1
        for col in range(1, 51):
            ws.cell(row=row, column=col, value=rng.getrandbits(128))
        if row % 500 == 0:
            wb.save(str(path))
            if path.stat().st_size >= _THRESHOLD_BYTES:
                break
        if row > 50000:
            break
    return path




@pytest.fixture(scope="session")
def secret_file(tmp_path_factory):
    """Write test HMAC secret to a temp file."""
    path = tmp_path_factory.mktemp("secrets") / "test_secret"
    path.write_text(TEST_SECRET)
    return path




@pytest.fixture(scope="session")
def _server_socket(tmp_path_factory):
    """Temp directory + socket path for the test gRPC server."""
    sock_dir = tmp_path_factory.mktemp("sidecar-sock")
    return str(sock_dir / "test.sock")


@pytest.fixture(scope="session")
def grpc_channel(_server_socket):
    server = _start_grpc_server(_server_socket)
    channel = grpc.insecure_channel(f"unix://{_server_socket}", options=_GRPC_OPTIONS)
    grpc.channel_ready_future(channel).result(timeout=5)
    yield channel
    channel.close()
    server.stop(grace=0)


@pytest.fixture(scope="session")
def grpc_stub(grpc_channel):
    """SidecarServiceStub connected to the no-auth server."""
    return sidecar_pb2_grpc.SidecarServiceStub(grpc_channel)




@pytest.fixture(scope="session")
def _auth_socket(tmp_path_factory):
    """Separate socket path for the authenticated server."""
    sock_dir = tmp_path_factory.mktemp("sidecar-auth")
    return str(sock_dir / "auth.sock")


@pytest.fixture(scope="session")
def grpc_channel_auth(_auth_socket):
    interceptor = _SyncTokenInterceptor(TEST_SECRET)
    server = _start_grpc_server(_auth_socket, interceptors=[interceptor])
    channel = grpc.insecure_channel(f"unix://{_auth_socket}", options=_GRPC_OPTIONS)
    grpc.channel_ready_future(channel).result(timeout=5)
    yield channel
    channel.close()
    server.stop(grace=0)


@pytest.fixture(scope="session")
def grpc_stub_auth(grpc_channel_auth):
    """SidecarServiceStub connected to the authenticated server."""
    return sidecar_pb2_grpc.SidecarServiceStub(grpc_channel_auth)
